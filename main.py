from pptx import Presentation
from pptx.util import Inches, Pt
from wand.image import Image
from wand.compat import nested
import os
import datetime

# App Intro
print("Welcome to pic2ppt")


# Function to create PPT with images+watermark
def create_ppt_file(directory, logo):
    # Validations
    if directory == "" or logo == "":
        return print("Invalid Input")

    # Creating New PPT
    new_pre = Presentation()
    slide_layout = new_pre.slide_layouts[6]

    files = []

    try:
        files = os.listdir(directory)
    except OSError as error:
        print(error)
        print("Path not found")
        return

    # Get Images
    for filename in files:  # return list
        if filename.endswith(".jpg"):
            image_path = directory + "/" + filename
            print(image_path)
            new_slide = new_pre.slides.add_slide(slide_layout)

            # Adding Heading and Subheading

            text_box = new_slide.shapes.add_textbox(left=Inches(1), top=Inches(0), width=Inches(3), height=Inches(0.5))

            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = "This is Heading"
            p.font.bold = True
            p.font.size = Pt(20)

            p = tf.add_paragraph()
            p.text = "And this is subheading"

            images = [
                logo,
                image_path
            ]
            with nested(Image(filename=images[0]), Image(filename=images[1])) as (watermark, canvas):
                watermark.transform(resize='x495')
                canvas.composite_channel("all_channels", watermark, "dissolve", 60, 60)

                left = top = Inches(1)

                # temp image path
                temp_img = directory + "/temp.jpg"

                canvas.save(filename=temp_img)
                print("Image Processed & Added To PPT")
                new_slide.shapes.add_picture(temp_img, left, top, height=Inches(5))

    # Save PPT
    print("Generating PPT ⏳ >>>>>>>>>>>>>>>>>>>>>>>>>")
    new_ppt = "New Presentation " + str(datetime.date.today()) + ".pptx"
    new_pre.save(directory + "/" + new_ppt)
    print("PPT File Generated Successfully")

    # Remove temp image file
    if os.path.exists(temp_img):
        os.remove(temp_img)
    else:
        print("The file does not exist")

    return "PPT File Generated Successfully"


# /function end

# inputs
print("Enter Image Directory Path")
image_dir = str(input())
print("Enter Logo Path")
logo_path = str(input())

# Call Function with dir and logo path
create_ppt_file(image_dir, logo_path)
